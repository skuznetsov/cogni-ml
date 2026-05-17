from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

OUT = '/Users/sergey/Projects/Crystal/cogni-ml/docs/talks/qwen35_inference_engine_30min.pptx'

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BG = RGBColor(13, 24, 34)
PANEL = RGBColor(24, 42, 54)
PANEL2 = RGBColor(31, 55, 68)
TEXT = RGBColor(238, 241, 232)
MUTED = RGBColor(166, 179, 174)
ACCENT = RGBColor(91, 214, 160)
ORANGE = RGBColor(240, 159, 83)
BLUE = RGBColor(105, 180, 255)
RED = RGBColor(245, 111, 111)

slides = [
    ("How We Built a Qwen 3.5/3.6 Inference Engine", "Native Crystal + Metal runtime for Qwen GGUF models", [
        "Goal: understand, control, and optimize inference",
        "Reference point: llama.cpp",
        "Method: implement → measure → refute → keep what survives",
    ], "title"),
    ("The Simple Mental Model", "LLM inference is repeated matrix math plus state updates", [
        "Text becomes tokens; tokens become vectors",
        "Each layer transforms vectors and updates memory",
        "The final vector predicts the next token",
        "Repeat until generation is complete",
    ], "flow"),
    ("Why Qwen 3.5/3.6 Was Hard", "It is a hybrid model, not just a plain transformer", [
        "Full-attention layers use KV cache",
        "Recurrent DeltaNet-style layers keep compact state",
        "Large FFN gate/up/down projections dominate traffic",
        "GGUF quantization: Q4_K, Q5_K, Q6_K, Q8_0",
    ], "two_col"),
    ("Engine Architecture", "Crystal controls the graph; Metal runs the hot math", [
        "CLI, probes, benchmark harnesses",
        "CPU/reference path for truth and debugging",
        "Metal kernels for matmul, attention, DeltaNet, top1, state copy",
    ], "stack"),
    ("Correctness Before Speed", "One wrong token can change the whole continuation", [
        "Layer-level and vector checks",
        "Logit/top1 comparisons",
        "Greedy token parity",
        "Approximate routes stay opt-in or proposal-only",
        "Exact verifier decides final tokens",
    ], "check"),
    ("Prefill ≠ Decode", "The same model has two very different workloads", [
        "Prefill: many known prompt tokens at once",
        "Decode: one new token at a time",
        "Prefill wants batch GEMM and chunking",
        "Decode wants low latency and fewer synchronizations",
    ], "split"),
    ("Prefill: What Worked", "Move from token-by-token to GPU-resident chunks", [
        "Chunked recurrent prefill on Metal",
        "Q4/Q5/Q6 batch GEMM for prompt tokens",
        "Skip unnecessary final-output work",
        "Keep consecutive recurrent layers on GPU",
        "Use attribution to find hot tensor shapes",
    ], "metrics_prefill"),
    ("Decode: What Worked", "Hide overhead and avoid unnecessary work per token", [
        "Wave scheduling groups layer work",
        "Fused greedy top1 avoids full logits when possible",
        "Small exact fusions remove extra kernels",
        "State preparation avoids first-touch latency",
        "N-gram/speculative paths help on matching workloads",
    ], "metrics_decode"),
    ("Speculative Decoding", "Draft fast, verify exactly", [
        "Draft proposes several tokens",
        "Exact target verifies the chunk",
        "Accept full prefix when it matches",
        "On mismatch: correct one token and continue",
    ], "spec"),
    ("Mathematical Shortcuts", "Do less proposal work, but keep exact verification", [
        "Projected-K / low-rank DeltaNet state",
        "PCA-updown and block surrogate ideas",
        "MTP / multi-token prediction experiments",
        "Raw-Q8 CUDA/DP4A as proposal-only until verified",
    ], "research"),
    ("The Workflow Was The Product", "Structured skepticism kept us from chasing noise", [
        "Landmark log: wins, failures, assumptions",
        "Paired A/B instead of one-off timings",
        "Compare against llama.cpp on same workload",
        "Record refutations to avoid loops",
        "Optimize after attribution names the wall",
    ], "process"),
    ("Useful Refutations", "A failed experiment is progress if it narrows the search", [
        "Microbench wins sometimes regressed full prefill",
        "Some fusions saved dispatches but hurt occupancy",
        "CPU draft lost to sync + body cost",
        "Fast full-row verifier could flip close logits",
        "GPU n-gram lookup was not the wall yet",
    ], "refute"),
    ("Where We Stand", "Mixed status, honest claims", [
        "llama.cpp remains a strong baseline",
        "Native decode has matched or beaten local snapshots",
        "First-run prefill is close and measurement-sensitive",
        "Prompt-cache and n-gram routes can win on right workloads",
        "vLLM/MLX teach scheduler and graph-boundary lessons",
    ], "status"),
    ("Next Frontier", "Turn observability into the next speed step", [
        "Quiet paired operator timing vs llama.cpp Metal",
        "Focus hot Q4_K 4096×12288 and recurrent Q5/Q6 shapes",
        "Continue exact n-gram/router policy work",
        "Keep approximate routes proposal-only until verified",
    ], "next"),
    ("Takeaways", "Five lessons from building the engine", [
        "Correctness gates first; speed second",
        "Prefill and decode are different products",
        "Data movement often dominates arithmetic",
        "Speculation helps only when economics line up",
        "A refutation log is an engineering asset",
    ], "takeaways"),
]


def add_bg(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    # subtle accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_textbox(slide, text, x, y, w, h, size=24, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.name = 'Aptos Display'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align: p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, 0.75, 0.45, 14.5, 0.65, size=34, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 0.78, 1.15, 13.2, 0.4, size=16, color=MUTED)


def add_bullets(slide, bullets, x=1.0, y=2.0, w=7.0, h=5.6, size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        p.font.name = 'Aptos'
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
    return box


def panel(slide, x, y, w, h, color=PANEL, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(54, 78, 91)
    shape.line.width = Pt(1)
    return shape


def footer(slide, idx):
    add_textbox(slide, f"Qwen 3.5/3.6 inference engine · {idx:02d}", 11.7, 8.45, 3.6, 0.25, size=9, color=RGBColor(117,132,130), align=PP_ALIGN.RIGHT)


def draw_flow(slide):
    labels = ["Text", "Tokens", "Vectors", "Layers + State", "Next token"]
    xs = [1.0, 3.6, 6.2, 8.8, 12.0]
    for i, (x, lab) in enumerate(zip(xs, labels)):
        panel(slide, x, 3.0, 2.0 if i != 3 else 2.4, 1.15, PANEL2, True)
        add_textbox(slide, lab, x+0.15, 3.37, 1.9 if i != 3 else 2.1, 0.3, size=18, bold=True, align=PP_ALIGN.CENTER)
        if i < len(xs)-1:
            con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + (2.0 if i != 3 else 2.4)), Inches(3.58), Inches(xs[i+1]), Inches(3.58))
            con.line.color.rgb = ACCENT
            con.line.width = Pt(2)


def draw_stack(slide):
    rows = [("CLI / probes / benchmarks", BLUE), ("Crystal reference + scheduler", ACCENT), ("Metal kernels", ORANGE)]
    y = 2.3
    for label, color in rows:
        panel(slide, 3.2, y, 9.5, 1.05, PANEL2, True)
        add_textbox(slide, label, 3.45, y+0.32, 9.0, 0.35, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        y += 1.35


def draw_spec(slide):
    nodes = [("Draft", 1.0, 3.0, BLUE), ("Candidates", 4.0, 3.0, ORANGE), ("Exact verifier", 7.3, 3.0, ACCENT), ("Accept / Correct", 11.2, 3.0, TEXT)]
    for label, x, y, color in nodes:
        panel(slide, x, y, 2.5, 1.15, PANEL2, True)
        add_textbox(slide, label, x+0.12, y+0.37, 2.25, 0.3, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    for i in range(len(nodes)-1):
        x1 = nodes[i][1] + 2.5; x2 = nodes[i+1][1]
        con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(3.58), Inches(x2), Inches(3.58))
        con.line.color.rgb = ACCENT; con.line.width = Pt(2)
    add_textbox(slide, "Final output remains exact because the target model verifies the accepted tokens.", 2.1, 5.0, 11.8, 0.5, size=20, color=MUTED, align=PP_ALIGN.CENTER)


def draw_metrics(slide, kind):
    if kind == 'prefill':
        items = [("~52.9 tok/s", "early pp64"), ("~358 tok/s", "after chunking + batching"), ("~450 tok/s", "prepared-state checkpoint")]
    else:
        items = [("+4–5%", "local decode snapshots vs llama.cpp"), ("~9–14 ms/tok", "repeat-heavy n-gram/spec routes"), ("exact", "verifier-gated output")]
    x = 8.7
    y = 2.05
    for value, label in items:
        panel(slide, x, y, 5.5, 1.25, PANEL2, True)
        add_textbox(slide, value, x+0.25, y+0.18, 5.0, 0.45, size=28, bold=True, color=ACCENT)
        add_textbox(slide, label, x+0.25, y+0.78, 5.0, 0.28, size=13, color=MUTED)
        y += 1.55


for idx, (title, subtitle, bullets, kind) in enumerate(slides, 1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    if kind == 'title':
        add_textbox(slide, "30-minute technical story · simple level", 0.9, 2.2, 8.0, 0.5, size=24, color=ACCENT, bold=True)
        add_bullets(slide, bullets, 0.95, 3.1, 7.5, 3.4, 24)
        panel(slide, 9.7, 2.2, 4.6, 3.9, PANEL2, True)
        add_textbox(slide, "Build → Measure → Refute", 10.0, 3.1, 4.0, 0.4, size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_textbox(slide, "Keep only what survives", 10.0, 3.75, 4.0, 0.35, size=18, color=MUTED, align=PP_ALIGN.CENTER)
    elif kind == 'flow':
        add_bullets(slide, bullets, 1.0, 5.1, 13.5, 2.2, 20)
        draw_flow(slide)
    elif kind == 'stack':
        add_bullets(slide, bullets, 1.0, 6.65, 13.5, 1.2, 18)
        draw_stack(slide)
    elif kind == 'spec':
        add_bullets(slide, bullets, 1.0, 5.75, 13.2, 1.6, 19)
        draw_spec(slide)
    elif kind == 'metrics_prefill':
        add_bullets(slide, bullets, 1.0, 2.05, 7.0, 5.5, 20)
        draw_metrics(slide, 'prefill')
    elif kind == 'metrics_decode':
        add_bullets(slide, bullets, 1.0, 2.05, 7.0, 5.5, 20)
        draw_metrics(slide, 'decode')
    elif kind == 'split':
        panel(slide, 1.0, 2.35, 6.5, 3.4, PANEL2, True)
        add_textbox(slide, "Prefill", 1.35, 2.75, 5.8, 0.4, size=28, bold=True, color=BLUE)
        add_textbox(slide, "Many known prompt tokens\nBatch GEMM + chunking\nThroughput-oriented", 1.35, 3.5, 5.6, 1.4, size=21, color=TEXT)
        panel(slide, 8.5, 2.35, 6.5, 3.4, PANEL2, True)
        add_textbox(slide, "Decode", 8.85, 2.75, 5.8, 0.4, size=28, bold=True, color=ORANGE)
        add_textbox(slide, "One new token at a time\nGEMV + low latency\nSynchronization-sensitive", 8.85, 3.5, 5.6, 1.4, size=21, color=TEXT)
        add_bullets(slide, bullets[-2:], 2.0, 6.3, 12.0, 1.2, 19)
    elif kind == 'takeaways':
        y = 2.0
        for i, b in enumerate(bullets, 1):
            add_textbox(slide, str(i), 1.15, y+0.05, 0.4, 0.35, size=20, bold=True, color=ACCENT)
            panel(slide, 1.75, y, 12.6, 0.7, PANEL2, True)
            add_textbox(slide, b, 2.05, y+0.18, 11.9, 0.25, size=18, bold=True)
            y += 0.95
    else:
        add_bullets(slide, bullets, 1.05, 2.0, 13.5, 5.6, 22 if len(bullets) <= 4 else 20)
    footer(slide, idx)

prs.save(OUT)
print(OUT)
